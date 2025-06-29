import os
import io
import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
import mimetypes
import logging

# Core libraries
import pytesseract
from PIL import Image
import pandas as pd

# Document processing libraries
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
import mammoth
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown2

# OCR libraries
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

# Language detection
try:
    from langdetect import detect, LangDetectError
    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False

# File type detection
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False

# Setup logging
logger = logging.getLogger(__name__)

@dataclass
class DocumentMetadata:
    """Metadata extracted from documents"""
    filename: str
    file_type: str
    mime_type: str
    file_size: int
    language: Optional[str] = None
    page_count: Optional[int] = None
    creation_date: Optional[datetime] = None
    modification_date: Optional[datetime] = None
    author: Optional[str] = None
    title: Optional[str] = None
    subject: Optional[str] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None

@dataclass
class ExtractedContent:
    """Content extracted from documents"""
    text: str
    metadata: DocumentMetadata
    tables: List[Dict] = None
    images: List[Dict] = None
    raw_content: Optional[str] = None
    ocr_confidence: Optional[float] = None
    processing_method: Optional[str] = None
    errors: List[str] = None

class DocumentTypeDetector:
    """Detect document type and MIME type"""
    
    SUPPORTED_EXTENSIONS = {
        # Text documents
        '.pdf': 'application/pdf',
        '.doc': 'application/msword',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.txt': 'text/plain',
        '.rtf': 'application/rtf',
        
        # Presentations
        '.ppt': 'application/vnd.ms-powerpoint',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        
        # Spreadsheets
        '.xls': 'application/vnd.ms-excel',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.csv': 'text/csv',
        
        # Images
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.png': 'image/png',
        '.tiff': 'image/tiff',
        '.tif': 'image/tiff',
        '.bmp': 'image/bmp',
        '.gif': 'image/gif',
        
        # Web formats
        '.html': 'text/html',
        '.htm': 'text/html',
        '.md': 'text/markdown',
        
        # Data formats
        '.json': 'application/json',
        '.xml': 'application/xml',
    }
    
    @classmethod
    def detect_file_type(cls, filepath: str) -> Tuple[str, str]:
        """Detect file type and MIME type"""
        file_ext = Path(filepath).suffix.lower()
        
        # First try extension-based detection
        if file_ext in cls.SUPPORTED_EXTENSIONS:
            mime_type = cls.SUPPORTED_EXTENSIONS[file_ext]
            return file_ext, mime_type
        
        # Fallback to python-magic for MIME type detection if available
        if MAGIC_AVAILABLE:
            try:
                mime_type = magic.from_file(filepath, mime=True)
                return file_ext, mime_type
            except Exception:
                pass
        
        # Final fallback to mimetypes
        mime_type, _ = mimetypes.guess_type(filepath)
        return file_ext, mime_type or 'application/octet-stream'
    
    @classmethod
    def is_supported(cls, filepath: str) -> bool:
        """Check if file type is supported"""
        file_ext, _ = cls.detect_file_type(filepath)
        return file_ext in cls.SUPPORTED_EXTENSIONS

class LanguageDetector:
    """Detect document language"""
    
    @staticmethod
    def detect_language(text: str) -> Optional[str]:
        """Detect language from text"""
        if not LANGDETECT_AVAILABLE or not text or len(text.strip()) < 10:
            return None
        
        try:
            # Clean text for better detection
            clean_text = ' '.join(text.split()[:100])  # Use first 100 words
            language = detect(clean_text)
            return language
        except LangDetectError:
            return None

class BaseExtractor:
    """Base class for document extractors"""
    
    def __init__(self):
        self.supported_extensions = []
        self.extractor_name = "base"
    
    def can_extract(self, file_path: str) -> bool:
        """Check if this extractor can handle the file"""
        ext = Path(file_path).suffix.lower()
        return ext in self.supported_extensions
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract content from file"""
        raise NotImplementedError
    
    def _get_basic_metadata(self, file_path: str, filename: str) -> DocumentMetadata:
        """Get basic file metadata"""
        file_stats = os.stat(file_path)
        file_ext, mime_type = DocumentTypeDetector.detect_file_type(file_path)
        
        return DocumentMetadata(
            filename=filename,
            file_type=file_ext,
            mime_type=mime_type,
            file_size=file_stats.st_size,
            creation_date=datetime.fromtimestamp(file_stats.st_ctime),
            modification_date=datetime.fromtimestamp(file_stats.st_mtime)
        )

class PDFExtractor(BaseExtractor):
    """Extract text from PDF files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pdf']
        self.extractor_name = "pdf"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from PDF"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        tables = []
        processing_method = None
        errors = []
        ocr_confidence = None
        
        try:
            # Try PyMuPDF first for text extraction
            doc = fitz.open(file_path)
            metadata.page_count = doc.page_count
            
            # Extract document metadata
            pdf_metadata = doc.metadata
            if pdf_metadata:
                metadata.title = pdf_metadata.get('title')
                metadata.author = pdf_metadata.get('author')
                metadata.subject = pdf_metadata.get('subject')
            
            # Extract text from each page
            text_parts = []
            has_text = False
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                page_text = page.get_text()
                
                if page_text.strip():
                    has_text = True
                    text_parts.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}")
            
            if has_text:
                text = "\n\n".join(text_parts)
                processing_method = "text_extraction"
            else:
                # PDF appears to be image-based, use OCR
                text, ocr_confidence = self._ocr_pdf(file_path)
                processing_method = "ocr"
            
            doc.close()
            
            # Try to extract tables using pdfplumber
            try:
                tables = self._extract_tables_pdfplumber(file_path)
            except Exception as e:
                errors.append(f"Table extraction failed: {str(e)}")
            
        except Exception as e:
            errors.append(f"PDF processing failed: {str(e)}")
            # Fallback to OCR
            try:
                text, ocr_confidence = self._ocr_pdf(file_path)
                processing_method = "ocr_fallback"
            except Exception as ocr_e:
                errors.append(f"OCR fallback failed: {str(ocr_e)}")
                text = f"Error: Could not extract text from PDF. {str(e)}"
        
        # Detect language and update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            tables=tables,
            processing_method=processing_method,
            ocr_confidence=ocr_confidence,
            errors=errors if errors else None
        )
    
    def _ocr_pdf(self, file_path: str) -> Tuple[str, float]:
        """OCR PDF using pdf2image and pytesseract"""
        try:
            from pdf2image import convert_from_path
            
            # Convert PDF to images
            images = convert_from_path(file_path, dpi=300)
            
            text_parts = []
            confidences = []
            
            for i, image in enumerate(images):
                # Get OCR data with confidence
                try:
                    ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    # Extract text and calculate confidence
                    page_text = pytesseract.image_to_string(image)
                    page_confidence = sum([int(conf) for conf in ocr_data['conf'] if int(conf) > 0])
                    page_confidence = page_confidence / len([c for c in ocr_data['conf'] if int(c) > 0]) if ocr_data['conf'] else 0
                    
                    if page_text.strip():
                        text_parts.append(f"--- Page {i+1} ---\n{page_text.strip()}")
                        confidences.append(page_confidence)
                except Exception:
                    # Fallback to simple OCR
                    page_text = pytesseract.image_to_string(image)
                    if page_text.strip():
                        text_parts.append(f"--- Page {i+1} ---\n{page_text.strip()}")
                        confidences.append(50.0)  # Default confidence
            
            text = "\n\n".join(text_parts)
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            return text, avg_confidence
            
        except ImportError:
            raise Exception("pdf2image not available. Install with: pip install pdf2image")
        except Exception as e:
            raise Exception(f"OCR processing failed: {str(e)}")
    
    def _extract_tables_pdfplumber(self, file_path: str) -> List[Dict]:
        """Extract tables using pdfplumber"""
        tables = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_tables = page.extract_tables()
                    
                    for table_num, table in enumerate(page_tables):
                        if table and len(table) > 1:
                            # Convert table to DataFrame for better structure
                            try:
                                df = pd.DataFrame(table[1:], columns=table[0] if table[0] else None)
                                
                                tables.append({
                                    'page': page_num + 1,
                                    'table_number': table_num + 1,
                                    'data': df.to_dict('records'),
                                    'headers': df.columns.tolist() if df.columns is not None else [],
                                    'rows': len(df),
                                    'columns': len(df.columns) if df.columns is not None else 0
                                })
                            except Exception:
                                # Fallback to raw table data
                                tables.append({
                                    'page': page_num + 1,
                                    'table_number': table_num + 1,
                                    'data': table,
                                    'headers': table[0] if table else [],
                                    'rows': len(table),
                                    'columns': len(table[0]) if table else 0
                                })
        except Exception:
            pass  # Tables extraction is optional
        
        return tables

class DOCXExtractor(BaseExtractor):
    """Extract text from DOCX files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.docx']
        self.extractor_name = "docx"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from DOCX"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        tables = []
        errors = []
        
        try:
            doc = Document(file_path)
            
            # Extract document properties
            try:
                props = doc.core_properties
                metadata.title = props.title
                metadata.author = props.author
                metadata.subject = props.subject
                metadata.creation_date = props.created
                metadata.modification_date = props.modified
            except Exception:
                pass  # Properties are optional
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            
            text = '\n\n'.join(paragraphs)
            
            # Extract tables
            for table_num, table in enumerate(doc.tables):
                table_data = []
                headers = []
                
                for row_num, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    
                    if row_num == 0:
                        headers = row_data
                    else:
                        table_data.append(row_data)
                
                if table_data:
                    try:
                        df = pd.DataFrame(table_data, columns=headers if headers else None)
                        tables.append({
                            'table_number': table_num + 1,
                            'data': df.to_dict('records'),
                            'headers': headers,
                            'rows': len(df),
                            'columns': len(df.columns) if df.columns is not None else len(headers)
                        })
                    except Exception:
                        # Fallback to raw table data
                        tables.append({
                            'table_number': table_num + 1,
                            'data': table_data,
                            'headers': headers,
                            'rows': len(table_data),
                            'columns': len(headers)
                        })
            
        except Exception as e:
            errors.append(f"DOCX processing failed: {str(e)}")
            text = f"Error: Could not extract text from DOCX. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            tables=tables,
            processing_method="docx_extraction",
            errors=errors if errors else None
        )

class DOCExtractor(BaseExtractor):
    """Extract text from DOC files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.doc']
        self.extractor_name = "doc"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from DOC using mammoth"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.extract_raw_text(docx_file)
                text = result.value
                
                if result.messages:
                    for message in result.messages:
                        if message.type == "warning" or message.type == "error":
                            errors.append(str(message))
            
        except Exception as e:
            errors.append(f"DOC processing failed: {str(e)}")
            text = f"Error: Could not extract text from DOC. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            processing_method="doc_extraction",
            errors=errors if errors else None
        )

class PPTXExtractor(BaseExtractor):
    """Extract text from PPTX files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx']
        self.extractor_name = "pptx"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from PPTX"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        
        try:
            prs = Presentation(file_path)
            
            # Extract presentation properties
            try:
                props = prs.core_properties
                metadata.title = props.title
                metadata.author = props.author
                metadata.subject = props.subject
                metadata.creation_date = props.created
                metadata.modification_date = props.modified
            except Exception:
                pass  # Properties are optional
            
            # Extract text from slides
            slide_texts = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text_parts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                
                if slide_text_parts:
                    slide_text = '\n'.join(slide_text_parts)
                    slide_texts.append(f"--- Slide {slide_num + 1} ---\n{slide_text}")
            
            text = '\n\n'.join(slide_texts)
            metadata.page_count = len(prs.slides)
            
        except Exception as e:
            errors.append(f"PPTX processing failed: {str(e)}")
            text = f"Error: Could not extract text from PPTX. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            processing_method="pptx_extraction",
            errors=errors if errors else None
        )

class ExcelExtractor(BaseExtractor):
    """Extract text and tables from Excel files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.xlsx', '.xls', '.csv']
        self.extractor_name = "excel"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text and tables from Excel"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        tables = []
        errors = []
        
        try:
            if file_path.endswith('.csv'):
                # Handle CSV files
                df = pd.read_csv(file_path)
                sheet_name = "Sheet1"
                sheets = {sheet_name: df}
            else:
                # Handle Excel files
                sheets = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            
            text_parts = []
            
            for sheet_name, df in sheets.items():
                if df.empty:
                    continue
                
                # Convert DataFrame to text representation
                sheet_text = f"--- Sheet: {sheet_name} ---\n"
                sheet_text += df.to_string(index=False)
                text_parts.append(sheet_text)
                
                # Store table data
                tables.append({
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'headers': df.columns.tolist(),
                    'rows': len(df),
                    'columns': len(df.columns)
                })
            
            text = '\n\n'.join(text_parts)
            
        except Exception as e:
            errors.append(f"Excel processing failed: {str(e)}")
            text = f"Error: Could not extract data from Excel file. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            tables=tables,
            processing_method="excel_extraction",
            errors=errors if errors else None
        )

class TextExtractor(BaseExtractor):
    """Extract text from plain text files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.txt', '.rtf']
        self.extractor_name = "text"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from text files"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        
        try:
            if file_path.endswith('.rtf'):
                # Handle RTF files
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(file_path, 'r', encoding='utf-8') as file:
                        rtf_content = file.read()
                        text = rtf_to_text(rtf_content)
                except ImportError:
                    errors.append("striprtf not available for RTF processing")
                    text = "Error: RTF processing requires striprtf package"
            else:
                # Handle plain text files with encoding detection
                encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as file:
                            text = file.read()
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    errors.append("Could not decode text file with common encodings")
                    text = "Error: Unable to decode text file"
            
        except Exception as e:
            errors.append(f"Text processing failed: {str(e)}")
            text = f"Error: Could not read text file. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            processing_method="text_extraction",
            errors=errors if errors else None
        )

class ImageExtractor(BaseExtractor):
    """Extract text from images using OCR"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif']
        self.extractor_name = "image_ocr"
        
        # Initialize EasyOCR if available
        self.easyocr_reader = None
        if EASYOCR_AVAILABLE:
            try:
                self.easyocr_reader = easyocr.Reader(['en'])  # Default to English
            except Exception:
                pass
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from images using OCR"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        ocr_confidence = 0
        errors = []
        processing_method = "tesseract_ocr"
        
        try:
            image = Image.open(file_path)
            
            # Get image metadata
            if hasattr(image, '_getexif') and image._getexif():
                exif = image._getexif()
                # Could extract more metadata from EXIF if needed
            
            # Try EasyOCR first if available
            if self.easyocr_reader:
                try:
                    results = self.easyocr_reader.readtext(file_path)
                    text_parts = []
                    confidences = []
                    
                    for (bbox, text_part, confidence) in results:
                        if confidence > 0.3:  # Filter low confidence results
                            text_parts.append(text_part)
                            confidences.append(confidence)
                    
                    text = ' '.join(text_parts)
                    ocr_confidence = sum(confidences) / len(confidences) if confidences else 0
                    processing_method = "easyocr"
                    
                except Exception as e:
                    errors.append(f"EasyOCR failed: {str(e)}")
            
            # Fallback to Tesseract
            if not text or len(text.strip()) < 10:
                try:
                    # Get OCR data with confidence
                    try:
                        ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                        
                        # Extract text
                        text = pytesseract.image_to_string(image)
                        
                        # Calculate average confidence
                        confidences = [int(conf) for conf in ocr_data['conf'] if int(conf) > 0]
                        ocr_confidence = sum(confidences) / len(confidences) if confidences else 0
                    except Exception:
                        # Fallback to simple OCR
                        text = pytesseract.image_to_string(image)
                        ocr_confidence = 50.0  # Default confidence
                    
                    processing_method = "tesseract_ocr"
                    
                except Exception as e:
                    errors.append(f"Tesseract OCR failed: {str(e)}")
                    text = f"Error: Could not extract text from image. {str(e)}"
            
            if not text.strip():
                text = "Warning: No text detected in image. The image may not contain readable text."
            
        except Exception as e:
            errors.append(f"Image processing failed: {str(e)}")
            text = f"Error: Could not process image. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            processing_method=processing_method,
            ocr_confidence=ocr_confidence,
            errors=errors if errors else None
        )

class HTMLExtractor(BaseExtractor):
    """Extract text from HTML files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.html', '.htm']
        self.extractor_name = "html"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from HTML"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        raw_content = ""
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
                raw_content = html_content
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract title
            title_tag = soup.find('title')
            if title_tag:
                metadata.title = title_tag.get_text().strip()
            
            # Extract meta information
            author_meta = soup.find('meta', attrs={'name': 'author'})
            if author_meta:
                metadata.author = author_meta.get('content')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text content
            text = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
        except Exception as e:
            errors.append(f"HTML processing failed: {str(e)}")
            text = f"Error: Could not extract text from HTML. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            raw_content=raw_content,
            processing_method="html_extraction",
            errors=errors if errors else None
        )

class MarkdownExtractor(BaseExtractor):
    """Extract text from Markdown files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.md', '.markdown']
        self.extractor_name = "markdown"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from Markdown"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        raw_content = ""
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                markdown_content = file.read()
                raw_content = markdown_content
            
            # Convert markdown to HTML then extract text
            html = markdown2.markdown(markdown_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
        except Exception as e:
            errors.append(f"Markdown processing failed: {str(e)}")
            text = f"Error: Could not extract text from Markdown. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            raw_content=raw_content,
            processing_method="markdown_extraction",
            errors=errors if errors else None
        )

class JSONExtractor(BaseExtractor):
    """Extract text from JSON files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.json']
        self.extractor_name = "json"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from JSON"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        raw_content = ""
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                json_data = json.load(file)
            
            # Extract text from JSON structure
            text_parts = self._extract_text_from_json(json_data)
            text = '\n'.join(text_parts)
            
            # Keep raw JSON as structured data
            raw_content = json.dumps(json_data, indent=2)
            
        except Exception as e:
            errors.append(f"JSON processing failed: {str(e)}")
            text = f"Error: Could not extract text from JSON. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            raw_content=raw_content,
            processing_method="json_extraction",
            errors=errors if errors else None
        )
    
    def _extract_text_from_json(self, obj, prefix="") -> List[str]:
        """Recursively extract text from JSON object"""
        text_parts = []
        
        if isinstance(obj, dict):
            for key, value in obj.items():
                current_prefix = f"{prefix}.{key}" if prefix else key
                if isinstance(value, (dict, list)):
                    text_parts.extend(self._extract_text_from_json(value, current_prefix))
                else:
                    text_parts.append(f"{current_prefix}: {str(value)}")
        
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                current_prefix = f"{prefix}[{i}]" if prefix else f"[{i}]"
                if isinstance(item, (dict, list)):
                    text_parts.extend(self._extract_text_from_json(item, current_prefix))
                else:
                    text_parts.append(f"{current_prefix}: {str(item)}")
        
        else:
            text_parts.append(str(obj))
        
        return text_parts

class XMLExtractor(BaseExtractor):
    """Extract text from XML files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.xml']
        self.extractor_name = "xml"
    
    def extract(self, file_path: str, filename: str) -> ExtractedContent:
        """Extract text from XML"""
        metadata = self._get_basic_metadata(file_path, filename)
        text = ""
        errors = []
        raw_content = ""
        
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # Extract text from XML elements
            text_parts = self._extract_text_from_xml(root)
            text = '\n'.join(text_parts)
            
            # Keep raw XML as structured data
            with open(file_path, 'r', encoding='utf-8') as file:
                raw_content = file.read()
            
        except Exception as e:
            errors.append(f"XML processing failed: {str(e)}")
            text = f"Error: Could not extract text from XML. {str(e)}"
        
        # Update metadata
        language = LanguageDetector.detect_language(text)
        metadata.language = language
        metadata.word_count = len(text.split()) if text else 0
        metadata.character_count = len(text) if text else 0
        
        return ExtractedContent(
            text=text,
            metadata=metadata,
            raw_content=raw_content,
            processing_method="xml_extraction",
            errors=errors if errors else None
        )
    
    def _extract_text_from_xml(self, element, prefix="") -> List[str]:
        """Recursively extract text from XML element"""
        text_parts = []
        
        current_prefix = f"{prefix}.{element.tag}" if prefix else element.tag
        
        # Add element text if it exists
        if element.text and element.text.strip():
            text_parts.append(f"{current_prefix}: {element.text.strip()}")
        
        # Add attributes
        if element.attrib:
            for attr, value in element.attrib.items():
                text_parts.append(f"{current_prefix}@{attr}: {value}")
        
        # Process child elements
        for child in element:
            text_parts.extend(self._extract_text_from_xml(child, current_prefix))
        
        return text_parts

class DocumentProcessor:
    """Main document processor that coordinates all extractors"""
    
    def __init__(self):
        self.extractors = [
            PDFExtractor(),
            DOCXExtractor(),
            DOCExtractor(),
            PPTXExtractor(),
            ExcelExtractor(),
            TextExtractor(),
            ImageExtractor(),
            HTMLExtractor(),
            MarkdownExtractor(),
            JSONExtractor(),
            XMLExtractor(),
        ]
        
        # Create a mapping of extensions to extractors for faster lookup
        self.extractor_map = {}
        for extractor in self.extractors:
            for ext in extractor.supported_extensions:
                self.extractor_map[ext] = extractor
    
    def get_supported_formats(self) -> List[str]:
        """Get list of all supported file formats"""
        formats = []
        for extractor in self.extractors:
            formats.extend(extractor.supported_extensions)
        return sorted(list(set(formats)))
    
    def is_supported(self, file_path: str) -> bool:
        """Check if file format is supported"""
        return DocumentTypeDetector.is_supported(file_path)
    
    def process_document(self, file_path: str, filename: str = None) -> ExtractedContent:
        """Process a document and extract its content"""
        if filename is None:
            filename = os.path.basename(file_path)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not self.is_supported(file_path):
            file_ext = Path(file_path).suffix.lower()
            supported = self.get_supported_formats()
            raise ValueError(f"Unsupported file type '{file_ext}'. Supported formats: {', '.join(supported)}")
        
        # Get the appropriate extractor
        file_ext = Path(file_path).suffix.lower()
        extractor = self.extractor_map.get(file_ext)
        
        if not extractor:
            raise ValueError(f"No extractor found for file type: {file_ext}")
        
        # Process the document
        try:
            result = extractor.extract(file_path, filename)
            logger.info(f"Successfully processed {filename} using {extractor.extractor_name}")
            return result
        except Exception as e:
            logger.error(f"Failed to process {filename}: {str(e)}")
            raise
    
    def process_multiple_documents(self, file_paths: List[str]) -> List[ExtractedContent]:
        """Process multiple documents"""
        results = []
        
        for file_path in file_paths:
            try:
                result = self.process_document(file_path)
                results.append(result)
            except Exception as e:
                # Create error result
                metadata = DocumentMetadata(
                    filename=os.path.basename(file_path),
                    file_type="unknown",
                    mime_type="unknown",
                    file_size=0
                )
                
                error_result = ExtractedContent(
                    text=f"Error processing file: {str(e)}",
                    metadata=metadata,
                    errors=[str(e)]
                )
                results.append(error_result)
        
        return results

# Global processor instance
document_processor = DocumentProcessor()

def serialize_extracted_content(content: ExtractedContent) -> dict:
    """Serialize ExtractedContent to JSON-compatible dict"""
    return {
        "text": content.text,
        "metadata": {
            "filename": content.metadata.filename,
            "file_type": content.metadata.file_type,
            "mime_type": content.metadata.mime_type,
            "file_size": content.metadata.file_size,
            "language": content.metadata.language,
            "page_count": content.metadata.page_count,
            "creation_date": content.metadata.creation_date.isoformat() if content.metadata.creation_date else None,
            "modification_date": content.metadata.modification_date.isoformat() if content.metadata.modification_date else None,
            "author": content.metadata.author,
            "title": content.metadata.title,
            "subject": content.metadata.subject,
            "word_count": content.metadata.word_count,
            "character_count": content.metadata.character_count,
        },
        "tables": content.tables,
        "images": content.images,
        "raw_content": content.raw_content,
        "ocr_confidence": content.ocr_confidence,
        "processing_method": content.processing_method,
        "errors": content.errors,
    } 