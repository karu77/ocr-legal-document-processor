import io
import os
import pytesseract
from PIL import Image
import tempfile
import subprocess
import sys
import re
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import spacy
from spacy.language import Language
from spacy.tokens import Doc
from translate import Translator
from langdetect import detect, LangDetectException
import pycountry

# Import for Word document processing
try:
    import docx2txt
    from docx import Document
    WORD_SUPPORT = True
except ImportError:
    WORD_SUPPORT = False

# Import for PDF processing
try:
    from pdf2image import convert_from_bytes, convert_from_path
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# Import for Excel processing
try:
    import openpyxl
    import xlrd
    import pandas as pd
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

# Import for PowerPoint processing
try:
    from pptx import Presentation
    POWERPOINT_SUPPORT = True
except ImportError:
    POWERPOINT_SUPPORT = False

# Import for OpenDocument processing
try:
    from odf.opendocument import load
    from odf.text import P
    from odf.element import Text
    from odf.teletype import extractText
    OPENDOCUMENT_SUPPORT = True
except ImportError:
    OPENDOCUMENT_SUPPORT = False

# Import for HTML processing
try:
    from bs4 import BeautifulSoup
    import lxml
    HTML_SUPPORT = True
except ImportError:
    HTML_SUPPORT = False

# Import for advanced text processing
try:
    import chardet
    import textract
    ADVANCED_TEXT_SUPPORT = True
except ImportError:
    ADVANCED_TEXT_SUPPORT = False

# Import for alternative OCR
try:
    import easyocr
    EASYOCR_SUPPORT = True
except ImportError:
    EASYOCR_SUPPORT = False

# Set up Google Cloud Vision client
# Ensure GOOGLE_APPLICATION_CREDENTIALS environment variable is set
# or provide credentials directly.
# For local development, it's common to set GOOGLE_APPLICATION_CREDENTIALS
# to the path of your service account key file.
# vision_client = vision.ImageAnnotatorClient()

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')

# Load spaCy model
try:
    nlp = spacy.load('en_core_web_sm')
except OSError:
    print("Downloading spaCy model...")
    subprocess.run([sys.executable, "-m", "spacy", "download", "en_core_web_sm"])
    nlp = spacy.load('en_core_web_sm')

def get_tesseract_lang_code(iso_code):
    """Convert ISO 639-1 (2-letter) to ISO 639-2 (3-letter) for Tesseract."""
    # Enhanced mapping for better language support including Hindi/Devanagari
    enhanced_mapping = {
        # Chinese variants
        'zh-cn': 'chi_sim',
        'zh-tw': 'chi_tra',
        'zh': 'chi_sim',  # Default Chinese to simplified
        
        # Indian languages
        'hi': 'hin',      # Hindi - Devanagari script
        'sa': 'san',      # Sanskrit - Devanagari script
        'mr': 'mar',      # Marathi - Devanagari script
        'ne': 'nep',      # Nepali - Devanagari script
        'bn': 'ben',      # Bengali
        'gu': 'guj',      # Gujarati
        'kn': 'kan',      # Kannada
        'ml': 'mal',      # Malayalam
        'or': 'ori',      # Odia
        'pa': 'pan',      # Punjabi
        'ta': 'tam',      # Tamil
        'te': 'tel',      # Telugu
        'ur': 'urd',      # Urdu
        
        # Other Asian languages
        'ja': 'jpn',      # Japanese
        'ko': 'kor',      # Korean
        'th': 'tha',      # Thai
        'vi': 'vie',      # Vietnamese
        
        # Middle Eastern
        'ar': 'ara',      # Arabic
        'fa': 'fas',      # Persian/Farsi
        'he': 'heb',      # Hebrew
        
        # European
        'en': 'eng',      # English
        'es': 'spa',      # Spanish
        'fr': 'fra',      # French
        'de': 'deu',      # German
        'it': 'ita',      # Italian
        'pt': 'por',      # Portuguese
        'ru': 'rus',      # Russian
        'pl': 'pol',      # Polish
        'nl': 'nld',      # Dutch
        'sv': 'swe',      # Swedish
        'no': 'nor',      # Norwegian
        'da': 'dan',      # Danish
        'fi': 'fin',      # Finnish
        'tr': 'tur',      # Turkish
        'el': 'ell',      # Greek
        'bg': 'bul',      # Bulgarian
        'cs': 'ces',      # Czech
        'sk': 'slk',      # Slovak
        'hr': 'hrv',      # Croatian
        'sr': 'srp',      # Serbian
        'sl': 'slv',      # Slovenian
        'et': 'est',      # Estonian
        'lv': 'lav',      # Latvian
        'lt': 'lit',      # Lithuanian
        'hu': 'hun',      # Hungarian
        'ro': 'ron',      # Romanian
        'uk': 'ukr',      # Ukrainian
    }
    
    # Check enhanced mapping first
    if iso_code in enhanced_mapping:
        return enhanced_mapping[iso_code]
    
    # Try pycountry as fallback
    try:
        lang = pycountry.languages.get(alpha_2=iso_code)
        if lang and hasattr(lang, 'alpha_3'):
            return lang.alpha_3
    except (AttributeError, KeyError):
        pass
    
    # Default to English if no mapping found
    return 'eng'

def perform_ocr_with_lang_detect(image_path_or_obj):
    """
    Performs OCR on an image, attempting to gracefully handle multiple languages,
    and translating to English if needed.
    """
    result = {
        'text': '',
        'detected_lang_name': 'English',
        'detected_lang_code': 'en',
        'original_text': '',
        'was_translated': False,
        'warning': None
    }

    try:
        # --- NEW LOGIC ---
        # 1. Perform a dual-language OCR pass with Tesseract first
        # This is surprisingly effective as Tesseract can handle scripts simultaneously.
        try:
            # Use English and Hindi packs together. Tesseract will pick the best fit.
            tesseract_dual_text = pytesseract.image_to_string(image_path_or_obj, lang='eng+hin')
        except pytesseract.TesseractError as e:
            # Handle cases where language packs might be missing
            print(f"Dual-language OCR failed, falling back to English. Error: {e}")
            tesseract_dual_text = pytesseract.image_to_string(image_path_or_obj, lang='eng')

        # 2. Use EasyOCR for a potentially better Hindi/mixed-language result.
        easyocr_text = None
        if EASYOCR_SUPPORT:
            try:
                # We use both 'hi' and 'en' to support mixed-language documents.
                easyocr_text = extract_text_with_easyocr(image_path_or_obj, detected_language='hi')
            except Exception as e:
                print(f"EasyOCR extraction failed during detection: {e}")
        
        # 3. Choose the best OCR result.
        # We check for Devanagari characters to determine if it's likely Hindi.
        devanagari_chars_regex = re.compile(r'[\u0900-\u097F]')
        
        final_ocr_text = tesseract_dual_text # Default to Tesseract's result
        iso_code = 'en' # Default to English

        # Check if EasyOCR provided a better result for Hindi
        if easyocr_text and devanagari_chars_regex.search(easyocr_text):
            # If EasyOCR text has more content and contains Hindi, prefer it.
            if len(easyocr_text) > len(tesseract_dual_text):
                 final_ocr_text = easyocr_text
                 iso_code = 'hi'

        # Check Tesseract's result if we haven't already decided on Hindi
        elif devanagari_chars_regex.search(tesseract_dual_text):
            iso_code = 'hi'
        
        # If we think it's Hindi, try to get a better name for it.
        if iso_code == 'hi':
            result['detected_lang_name'] = 'Hindi'
            result['detected_lang_code'] = 'hi'
        else:
            # If no Devanagari, we can try to detect other languages.
            try:
                iso_code = detect(final_ocr_text[:2000]) if final_ocr_text else 'en'
                lang_obj = pycountry.languages.get(alpha_2=iso_code)
                result['detected_lang_name'] = lang_obj.name if lang_obj else iso_code.upper()
                result['detected_lang_code'] = iso_code
            except Exception:
                # Fallback to English if detection fails
                result['detected_lang_name'] = 'English (assumed)'
                result['detected_lang_code'] = 'en'
                iso_code = 'en'
        
        if not final_ocr_text.strip():
            result['warning'] = "No text detected in the document."
            return result
            
        result['original_text'] = final_ocr_text

        # 4. SET FINAL TEXT (NO AUTO-TRANSLATION)
        # The extracted text is the original text. Translation will be a separate user action.
        result['text'] = final_ocr_text

        # Update warning if language was not English, to inform the user.
        if iso_code.lower() not in ['en', 'eng']:
            result['warning'] = f"Document processed in its original language ({result['detected_lang_name']}). Use the 'Translate' button to translate it."
        
        return result

    except Exception as e:
        print(f"OCR with language detection failed: {e}")
        result['text'] = f"Error: OCR processing failed - {str(e)}"
        return result

def auto_translate_to_english(text, source_language_code):
    """
    Automatically translate text to English using multiple translation services.
    Tries Hugging Face API first, then other services.
    """
    try:
        # Method 1: Try Hugging Face Translation API (if available)
        translated = translate_with_huggingface(text, source_language_code, 'en')
        if translated and translated != text:
            return translated
    except Exception as e:
        print(f"Hugging Face translation failed: {e}")
    
    try:
        # Method 2: Try MyMemory API
        translated = translate_with_mymemory(text, source_language_code, 'en')
        if translated and translated != text:
            return translated
    except Exception as e:
        print(f"MyMemory translation failed: {e}")
    
    try:
        # Method 3: Try Google Translate (if available)
        translated = translate_with_googletrans(text, source_language_code, 'en')
        if translated and translated != text:
            return translated
    except Exception as e:
        print(f"Google Translate failed: {e}")
    
    return text  # Return original if all translation methods fail

def translate_with_huggingface(text, source_lang, target_lang='en'):
    """
    Translate text using Hugging Face Inference API
    """
    try:
        import requests
        import os
        
        # Get Hugging Face API token from environment
        hf_token = os.getenv('HUGGINGFACE_API_TOKEN')
        if not hf_token:
            print("Hugging Face API token not found in environment variables")
            return {'success': False, 'translated_text': text}
        
        # Language mapping for Hugging Face models
        lang_mapping = {
            'hi': 'hi',  # Hindi
            'es': 'es',  # Spanish  
            'fr': 'fr',  # French
            'de': 'de',  # German
            'it': 'it',  # Italian
            'pt': 'pt',  # Portuguese
            'ru': 'ru',  # Russian
            'ja': 'ja',  # Japanese
            'ko': 'ko',  # Korean
            'zh': 'zh',  # Chinese
            'ar': 'ar',  # Arabic
        }
        
        source_lang_mapped = lang_mapping.get(source_lang, source_lang)
        
        # Use Helsinki-NLP translation models (very good quality)
        model_name = f"Helsinki-NLP/opus-mt-{source_lang_mapped}-en"
        
        headers = {
            "Authorization": f"Bearer {hf_token}",
            "Content-Type": "application/json"
        }
        
        # Split text into chunks for better processing
        chunks = split_text_for_translation(text, max_length=500)
        translated_chunks = []
        
        for chunk in chunks:
            if not chunk.strip():
                continue
                
            payload = {
                "inputs": chunk,
                "options": {"wait_for_model": True}
            }
            
            response = requests.post(
                f"https://api-inference.huggingface.co/models/{model_name}",
                headers=headers,
                json=payload,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                if isinstance(result, list) and len(result) > 0:
                    translated_text = result[0].get('translation_text', chunk)
                    translated_chunks.append(translated_text)
                else:
                    # In case of unexpected success response, append original
                    translated_chunks.append(chunk)
            else:
                print(f"Hugging Face API error: {response.status_code}")
                # On error, we don't have a valid translation for this chunk
                return {'success': False, 'translated_text': text}
                
        return {'success': True, 'translated_text': ' '.join(translated_chunks)}
        
    except Exception as e:
        print(f"Hugging Face translation error: {e}")
        return {'success': False, 'translated_text': text}

def translate_with_mymemory(text, source_lang, target_lang='en'):
    """
    Translate text using MyMemory API (free service)
    """
    try:
        import requests
        import urllib.parse
        
        chunks = split_text_for_translation(text, max_length=500)
        translated_chunks = []
        
        for chunk in chunks:
            if not chunk.strip():
                continue
                
            encoded_text = urllib.parse.quote(chunk)
            url = f"https://api.mymemory.translated.net/get?q={encoded_text}&langpair={source_lang}|{target_lang}"
            
            response = requests.get(url, timeout=10)
            
            if response.status_code == 200:
                data = response.json()
                if data.get('responseStatus') == 200:
                    translated_text = data['responseData']['translatedText']
                    translated_chunks.append(translated_text)
                else:
                    # MyMemory API reported an error
                    return {'success': False, 'translated_text': text}
            else:
                # HTTP error
                return {'success': False, 'translated_text': text}
                
        return {'success': True, 'translated_text': ' '.join(translated_chunks)}
        
    except Exception as e:
        print(f"MyMemory translation error: {e}")
        return {'success': False, 'translated_text': text}

def translate_with_googletrans(text, source_lang, target_lang='en'):
    """
    Translate text using googletrans library (if installed)
    """
    try:
        from googletrans import Translator
        
        translator = Translator()
        chunks = split_text_for_translation(text, max_length=5000)  # Google allows longer text
        translated_chunks = []
        
        for chunk in chunks:
            if not chunk.strip():
                continue
                
            result = translator.translate(chunk, src=source_lang, dest=target_lang)
            translated_chunks.append(result.text)
            
        return {'success': True, 'translated_text': ' '.join(translated_chunks)}
        
    except ImportError:
        print("googletrans library not installed. Install with: pip install googletrans==4.0.0-rc1")
        return {'success': False, 'translated_text': text}
    except Exception as e:
        print(f"Google Translate error: {e}")
        return {'success': False, 'translated_text': text}

def split_text_for_translation(text, max_length=500):
    """
    Split text into smaller chunks for translation APIs
    """
    if len(text) <= max_length:
        return [text]
    
    # Split by sentences first
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        if current_length + len(sentence) > max_length and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_length = len(sentence)
        else:
            current_chunk.append(sentence)
            current_length += len(sentence)
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def install_poppler_guide():
    """Return installation guide for Poppler on Windows"""
    return """
    To enable PDF processing, please install Poppler:
    
    Option 1 - Using Chocolatey (Run PowerShell as Administrator):
    choco install poppler
    
    Option 2 - Manual Installation:
    1. Download Poppler from: https://github.com/oschwartz10612/poppler-windows/releases/
    2. Extract to C:\\poppler
    3. Add C:\\poppler\\Library\\bin to your PATH environment variable
    4. Restart your command prompt/IDE
    
    Option 3 - Using conda:
    conda install -c conda-forge poppler
    """

def extract_text_from_excel(filepath):
    """Extract text from Excel files (XLS, XLSX)"""
    if not EXCEL_SUPPORT:
        return "Error: Excel support not installed. Please install: pip install openpyxl xlrd pandas"
    
    try:
        # Try to detect file format and read accordingly
        file_extension = os.path.splitext(filepath)[1].lower()
        
        if file_extension == '.xlsx':
            # Use openpyxl for .xlsx files
            workbook = openpyxl.load_workbook(filepath, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell))
                    if row_text:
                        text_parts.append(" | ".join(row_text))
                text_parts.append("")
            
            return "\n".join(text_parts)
        
        elif file_extension == '.xls':
            # Use pandas for .xls files (handles both xlrd and openpyxl)
            try:
                excel_file = pd.ExcelFile(filepath)
                text_parts = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(filepath, sheet_name=sheet_name)
                    text_parts.append(f"=== Sheet: {sheet_name} ===")
                    
                    # Convert DataFrame to text
                    for index, row in df.iterrows():
                        row_text = []
                        for value in row:
                            if pd.notna(value) and str(value).strip():
                                row_text.append(str(value))
                        if row_text:
                            text_parts.append(" | ".join(row_text))
                    text_parts.append("")
                
                return "\n".join(text_parts)
            except Exception as e:
                return f"Error reading XLS file: {e}"
        
    except Exception as e:
        return f"Error extracting text from Excel file: {e}"

def extract_text_from_powerpoint(filepath):
    """Extract text from PowerPoint files (PPT, PPTX)"""
    if not POWERPOINT_SUPPORT:
        return "Error: PowerPoint support not installed. Please install: pip install python-pptx"
    
    try:
        presentation = Presentation(filepath)
        text_parts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                    
                # Handle tables in slides
                if shape.shape_type == 19:  # Table
                    try:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))
                    except:
                        pass
            
            text_parts.append("")
        
        return "\n".join(text_parts)
    
    except Exception as e:
        return f"Error extracting text from PowerPoint: {e}"

def extract_text_from_opendocument(filepath):
    """Extract text from OpenDocument files (ODT, ODS, ODP)"""
    if not OPENDOCUMENT_SUPPORT:
        return "Error: OpenDocument support not installed. Please install: pip install odfpy"
    
    try:
        doc = load(filepath)
        text_parts = []
        
        # Extract all text elements
        for element in doc.getElementsByType(P):
            text = extractText(element)
            if text.strip():
                text_parts.append(text.strip())
        
        return "\n".join(text_parts) if text_parts else "No text found in OpenDocument file"
    
    except Exception as e:
        return f"Error extracting text from OpenDocument: {e}"

def extract_text_from_html(filepath):
    """Extract text from HTML files"""
    if not HTML_SUPPORT:
        return "Error: HTML support not installed. Please install: pip install beautifulsoup4 lxml"
    
    try:
        with open(filepath, 'r', encoding='utf-8') as file:
            content = file.read()
        
        soup = BeautifulSoup(content, 'lxml')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text if text else "No text found in HTML file"
    
    except Exception as e:
        return f"Error extracting text from HTML: {e}"

def extract_text_from_csv(filepath):
    """Extract text from CSV files"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                df = pd.read_csv(filepath, encoding=encoding)
                text_parts = []
                
                # Add column headers
                headers = " | ".join(str(col) for col in df.columns)
                text_parts.append(f"=== Headers ===\n{headers}\n")
                
                # Add data rows
                text_parts.append("=== Data ===")
                for index, row in df.iterrows():
                    row_text = []
                    for value in row:
                        if pd.notna(value) and str(value).strip():
                            row_text.append(str(value))
                    if row_text:
                        text_parts.append(" | ".join(row_text))
                
                return "\n".join(text_parts)
            
            except UnicodeDecodeError:
                continue
        
        return "Error: Unable to decode CSV file with common encodings"
    
    except Exception as e:
        return f"Error extracting text from CSV: {e}"

def extract_text_with_textract(filepath):
    """Fallback text extraction using textract for various formats"""
    if not ADVANCED_TEXT_SUPPORT:
        return "Error: Advanced text support not installed. Please install: pip install textract"
    
    try:
        text = textract.process(filepath, encoding='utf-8')
        if isinstance(text, bytes):
            text = text.decode('utf-8')
        return text.strip() if text else "No text could be extracted"
    
    except Exception as e:
        return f"Error extracting text with textract: {e}"

def detect_file_encoding(filepath):
    """Detect file encoding for better text extraction"""
    try:
        with open(filepath, 'rb') as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            return result.get('encoding', 'utf-8')
    except:
        return 'utf-8'

def extract_text_with_easyocr(filepath, detected_language='en'):
    """Alternative OCR using EasyOCR for better multilingual support, especially for Hindi/Devanagari"""
    if not EASYOCR_SUPPORT:
        return None
    
    try:
        # Configure languages based on detected language
        # EasyOCR uses different language codes than Tesseract
        easyocr_lang_mapping = {
            'hi': ['hi', 'en'],    # Hindi + English for mixed content
            'hin': ['hi', 'en'],   # Tesseract code to EasyOCR
            'bn': ['bn', 'en'],    # Bengali
            'ta': ['ta', 'en'],    # Tamil
            'te': ['te', 'en'],    # Telugu
            'kn': ['kn', 'en'],    # Kannada
            'ml': ['ml', 'en'],    # Malayalam
            'gu': ['gu', 'en'],    # Gujarati
            'or': ['or', 'en'],    # Odia
            'pa': ['pa', 'en'],    # Punjabi
            'ur': ['ur', 'en'],    # Urdu
            'ar': ['ar', 'en'],    # Arabic
            'fa': ['fa', 'en'],    # Persian
            'zh': ['ch_sim', 'en'], # Chinese Simplified
            'ja': ['ja', 'en'],    # Japanese
            'ko': ['ko', 'en'],    # Korean
            'th': ['th', 'en'],    # Thai
            'vi': ['vi', 'en'],    # Vietnamese
            'ru': ['ru', 'en'],    # Russian
            'fr': ['fr', 'en'],    # French
            'de': ['de', 'en'],    # German
            'es': ['es', 'en'],    # Spanish
            'pt': ['pt', 'en'],    # Portuguese
            'it': ['it', 'en'],    # Italian
        }
        
        # Get appropriate language list for EasyOCR
        languages = easyocr_lang_mapping.get(detected_language, ['en'])
        
        print(f"Using EasyOCR with languages: {languages}")
        reader = easyocr.Reader(languages)
        result = reader.readtext(filepath)
        
        # Extract text from results with lower confidence threshold for non-Latin scripts
        confidence_threshold = 0.3 if detected_language in ['hi', 'hin', 'bn', 'ta', 'te', 'kn', 'ml', 'gu', 'or', 'pa', 'ur', 'ar', 'fa'] else 0.5
        
        text_parts = []
        for (bbox, text, confidence) in result:
            if confidence > confidence_threshold:
                text_parts.append(text.strip())
        
        if text_parts:
            # For languages like Hindi, join words more naturally
            if detected_language in ['hi', 'hin']:
                return " ".join(text_parts)  # Join with spaces for better readability
            else:
                return "\n".join(text_parts)
        else:
            return None
    
    except Exception as e:
        print(f"EasyOCR failed: {e}")
        return None

def extract_text_from_docx(filepath):
    """Extract text from DOCX files"""
    try:
        # Try docx2txt first (simpler)
        text = docx2txt.process(filepath)
        if text and text.strip():
            return text.strip()
        
        # Fallback to python-docx
        doc = Document(filepath)
        text_parts = []
        for paragraph in doc.paragraphs:
            text_parts.append(paragraph.text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"Error extracting text from DOCX: {e}"

def extract_text_from_doc(filepath):
    """Extract text from DOC files using antiword or textract"""
    try:
        # Try using antiword command line tool
        result = subprocess.run(['antiword', filepath], 
                              capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout:
            return result.stdout
    except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.SubprocessError):
        pass
    
    # Fallback to textract
    return extract_text_with_textract(filepath)

def extract_text_from_txt(filepath):
    """Extract text from TXT files"""
    try:
        # First try to detect encoding
        if ADVANCED_TEXT_SUPPORT:
            encoding = detect_file_encoding(filepath)
            encodings = [encoding, 'utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        else:
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        return "Error: Unable to decode text file with common encodings"
    except Exception as e:
        return f"Error reading text file: {e}"

def setup_poppler_path():
    """Setup Poppler path from project directory"""
    from pathlib import Path
    
    # Get the project root directory
    current_file = Path(__file__)
    project_root = current_file.parent.parent.parent
    
    # Check for Poppler in project directory
    poppler_paths = [
        project_root / "poppler" / "poppler-23.11.0" / "Library" / "bin",
        project_root / "poppler" / "Library" / "bin",
        project_root / "poppler" / "bin"
    ]
    
    for poppler_path in poppler_paths:
        if poppler_path.exists() and (poppler_path / "pdftoppm.exe").exists():
            # Add to PATH if not already there
            poppler_str = str(poppler_path)
            current_path = os.environ.get('PATH', '')
            if poppler_str not in current_path:
                os.environ['PATH'] = poppler_str + os.pathsep + current_path
                print(f"Added Poppler to PATH: {poppler_str}")
            return True
    
    return False

def check_poppler_installation():
    """Check if Poppler is installed and accessible"""
    # First try to setup from project directory
    setup_poppler_path()
    
    try:
        # Try to run pdftoppm command
        result = subprocess.run(['pdftoppm', '-h'], 
                              capture_output=True, timeout=5)
        return result.returncode == 0
    except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.SubprocessError):
        return False

def process_pdf_fallback(filepath):
    """Fallback PDF processing using PyPDF2 for text extraction"""
    try:
        import PyPDF2
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---\n{text.strip()}")
            
            if text_parts:
                return "\n\n".join(text_parts)
            else:
                return "Error: PDF appears to be image-based. OCR processing requires Poppler installation."
    except ImportError:
        return "Error: PyPDF2 not installed. Please install: pip install PyPDF2"
    except Exception as e:
        return f"Error processing PDF: {e}"

def process_ocr(filepath, filename):
    """Process OCR for various document types with automatic translation to English."""
    file_extension = os.path.splitext(filename)[1].lower()
    
    # This dictionary will hold our final results
    ocr_result = {
        'text': '',
        'detected_lang_name': 'English',
        'detected_lang_code': 'en',
        'original_text': '',
        'was_translated': False,
        'warning': None
    }

    # Helper to update results with auto-translation support
    def update_result(data, is_ocr=False):
        if is_ocr:
            # For OCR results, the auto-translation is already handled
            ocr_result.update(data)
        else:
            # For non-OCR text, check if we need to translate
            ocr_result['text'] = data
            ocr_result['original_text'] = data
            
            # Try to detect language and auto-translate if not English
            if data and not data.startswith("Error:"):
                try:
                    from langdetect import detect
                    detected_lang = detect(data[:2000])  # Use first 2000 chars for detection
                    
                    if detected_lang.lower() not in ['en', 'eng']:
                        print(f"🔄 Auto-translating {detected_lang} document to English...")
                        translated = auto_translate_to_english(data, detected_lang)
                        
                        if translated and translated != data:
                            ocr_result['text'] = translated
                            ocr_result['was_translated'] = True
                            ocr_result['detected_lang_code'] = detected_lang
                            try:
                                lang_obj = pycountry.languages.get(alpha_2=detected_lang)
                                ocr_result['detected_lang_name'] = lang_obj.name if lang_obj else detected_lang.upper()
                            except:
                                ocr_result['detected_lang_name'] = detected_lang.upper()
                            ocr_result['warning'] = f"Document was automatically translated from {ocr_result['detected_lang_name']} to English."
                            print(f"✅ Successfully translated document to English")
                        else:
                            ocr_result['warning'] = f"Auto-translation from {detected_lang} failed. Showing original text."
                            
                except Exception as e:
                    print(f"Language detection/translation failed: {e}")
                    # Continue with original text

    # PDF Files
    if file_extension == '.pdf':
        if not PDF_SUPPORT:
            update_result("Error: pdf2image not installed. Please install: pip install pdf2image")
        elif not check_poppler_installation():
            print("Poppler not found, trying fallback method...")
            fallback_text = process_pdf_fallback(filepath)
            if "Error:" in fallback_text:
                fallback_text += f"\\n\\n{install_poppler_guide()}"
            update_result(fallback_text)
        else:
            try:
                print("📄 Processing PDF with auto-translation to English...")
                images = convert_from_path(filepath, dpi=300)
                all_texts = []
                original_texts = []
                detected_langs = []
                was_any_translated = False
                
                for i, image in enumerate(images):
                    page_result = perform_ocr_with_lang_detect(image)
                    all_texts.append(f"--- Page {i+1} ---\\n{page_result['text'].strip()}")
                    original_texts.append(f"--- Page {i+1} ---\\n{page_result.get('original_text', page_result['text']).strip()}")
                    
                    if page_result['detected_lang_code'] not in detected_langs:
                        detected_langs.append(page_result['detected_lang_code'])
                    
                    if page_result.get('was_translated', False):
                        was_any_translated = True
                
                # Consolidate results
                update_result({
                    'text': "\\n\\n".join(all_texts),
                    'original_text': "\\n\\n".join(original_texts),
                    'detected_lang_name': ', '.join(lang.capitalize() for lang in detected_langs),
                    'detected_lang_code': ','.join(detected_langs),
                    'was_translated': was_any_translated,
                    'warning': "Multi-page PDF was automatically processed and translated to English." if was_any_translated else None
                }, is_ocr=True)

            except Exception as e:
                print(f"Error processing PDF with Tesseract/Poppler: {e}")
                update_result(f"Error: Failed to extract text from PDF. {e}\\n\\n{install_poppler_guide()}")

    # Image Files (OCR Processing with auto-translation)
    elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.tif']:
        try:
            print(f"🖼️  Processing image with auto-translation to English: {filename}")
            update_result(perform_ocr_with_lang_detect(filepath), is_ocr=True)
            if not ocr_result['text'].strip() and not ocr_result['warning']:
                ocr_result['warning'] = "No text detected in image. The image may not contain readable text."
        except Exception as e:
            print(f"Error processing image with OCR: {e}")
            update_result(f"Error: Failed to extract text from image. {e}")

    # Microsoft Word Documents
    elif file_extension in ['.docx']:
        if not WORD_SUPPORT:
            update_result("Error: Word document support not installed. Please install: pip install python-docx docx2txt")
        else:
            update_result(extract_text_from_docx(filepath))

    elif file_extension in ['.doc']:
        update_result(extract_text_from_doc(filepath))

    # Other file types that don't need image-based OCR
    else:
        text_content = ""
        if file_extension in ['.xlsx', '.xls']:
            text_content = extract_text_from_excel(filepath)
        elif file_extension in ['.pptx', '.ppt']:
            text_content = extract_text_from_powerpoint(filepath)
        elif file_extension in ['.odt', '.ods', '.odp']:
            text_content = extract_text_from_opendocument(filepath)
        elif file_extension in ['.txt', '.text']:
            text_content = extract_text_from_txt(filepath)
        elif file_extension in ['.csv']:
            text_content = extract_text_from_csv(filepath)
        elif file_extension in ['.html', '.htm', '.xml']:
            text_content = extract_text_from_html(filepath)
        elif file_extension in ['.rtf']:
            try:
                import striprtf
                with open(filepath, 'r', encoding='utf-8') as file:
                    rtf_content = file.read()
                    text_content = striprtf.striprtf(rtf_content)
            except ImportError:
                text_content = "Error: RTF support not installed. Please install: pip install striprtf"
            except Exception as e:
                text_content = f"Error processing RTF file: {e}"
        else:
             text_content = f"Error: Unsupported file type '{file_extension}'."
        
        update_result(text_content)
    
    if not ocr_result['text'] and not ocr_result['warning']:
         ocr_result['text'] = "No text could be extracted from this document."

    return ocr_result

def clean_text(text: str) -> str:
    """Clean and format extracted text"""
    if not text:
        return ""
    
    # Basic cleanup
    text = re.sub(r'\s+', ' ', text)  # Remove multiple spaces
    text = re.sub(r'[^\x00-\x7F]+', '', text)  # Remove non-ASCII characters
    text = text.strip()
    
    # Process with spaCy
    doc = nlp(text)
    
    # Reconstruct text with proper formatting
    sentences = []
    current_paragraph = []
    
    for sent in doc.sents:
        # Clean the sentence
        clean_sent = re.sub(r'\s+', ' ', sent.text.strip())
        if clean_sent:
            current_paragraph.append(clean_sent)
            
            # Check for paragraph breaks (based on newlines or semantic breaks)
            if '\n' in sent.text or len(current_paragraph) > 5:
                sentences.append(' '.join(current_paragraph))
                current_paragraph = []
    
    # Add any remaining sentences
    if current_paragraph:
        sentences.append(' '.join(current_paragraph))
    
    # Join paragraphs with double newlines
    cleaned_text = '\n\n'.join(sentences)
    
    # Fix common OCR issues
    cleaned_text = re.sub(r'[lI](?=\d)', '1', cleaned_text)  # Fix common l/I to 1 confusion
    cleaned_text = re.sub(r'O(?=\d)', '0', cleaned_text)  # Fix common O to 0 confusion
    cleaned_text = re.sub(r'(?<=\d)[oO]', '0', cleaned_text)  # Fix common o/O to 0 confusion
    
    return cleaned_text

def summarize_text(text: str) -> str:
    """Generate a summary of the text using extractive summarization"""
    if not text:
        return ""
    
    # Process with spaCy
    doc = nlp(text)
    
    # Calculate sentence importance scores
    sentence_scores = {}
    for sent in doc.sents:
        # Score based on length and position
        score = len([token for token in sent if not token.is_stop and token.is_alpha])
        if sent.start == 0:  # First sentence gets a boost
            score *= 1.25
        sentence_scores[sent.text] = score
    
    # Get top sentences (about 30% of original)
    num_sentences = max(3, len(list(doc.sents)) // 3)
    summary_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:num_sentences]
    
    # Reconstruct summary in original order
    summary = []
    for sent in doc.sents:
        if sent.text in dict(summary_sentences):
            summary.append(sent.text.strip())
    
    return '\n\n'.join(summary)

def extract_key_points(text: str) -> list:
    """Extract key points from the text"""
    if not text:
        return []
    
    # Process with spaCy
    doc = nlp(text)
    
    # Extract important sentences based on various criteria
    key_points = []
    for sent in doc.sents:
        # Check for important features
        has_entity = any(ent.label_ in ['ORG', 'PERSON', 'DATE', 'MONEY', 'PERCENT'] for ent in sent.ents)
        has_numbers = any(token.like_num for token in sent)
        has_key_phrase = any(phrase in sent.text.lower() for phrase in [
            'must', 'shall', 'will', 'agree', 'require', 'important',
            'deadline', 'payment', 'terms', 'condition'
        ])
        
        if has_entity or has_numbers or has_key_phrase:
            clean_sent = re.sub(r'\s+', ' ', sent.text.strip())
            if clean_sent and clean_sent not in key_points:
                key_points.append(clean_sent)
    
    return key_points

def translate_text(text: str, target_language: str) -> str:
    """Translate text to target language"""
    if not text or not target_language:
        return ""
    
    try:
        # Initialize translator
        translator = Translator(to_lang=target_language)
        
        # Split text into manageable chunks (API limits)
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sent_tokenize(text):
            # Add sentence length plus a space
            sent_length = len(sentence) + 1
            
            if current_length + sent_length > 500:  # API typically has limits
                chunks.append(' '.join(current_chunk))
                current_chunk = [sentence]
                current_length = sent_length
            else:
                current_chunk.append(sentence)
                current_length += sent_length
        
        # Add any remaining chunk
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        # Translate each chunk
        translated_chunks = []
        for chunk in chunks:
            try:
                translated = translator.translate(chunk)
                if translated:
                    translated_chunks.append(translated)
            except Exception as e:
                print(f"Translation error for chunk: {e}")
                translated_chunks.append(chunk)  # Keep original on error
        
        # Join translated chunks
        return '\n\n'.join(translated_chunks)
        
    except Exception as e:
        print(f"Translation error: {e}")
        return text  # Return original text on error

def compare_documents(text1: str, text2: str) -> dict:
    """Compare two documents and return differences"""
    if not text1 or not text2:
        return {
            "similarity_percentage": 0,
            "differences": [],
            "common_points": []
        }
    
    # Process both texts
    doc1 = nlp(text1)
    doc2 = nlp(text2)
    
    # Compare sentences
    sentences1 = [sent.text.strip() for sent in doc1.sents]
    sentences2 = [sent.text.strip() for sent in doc2.sents]
    
    # Find common and different sentences
    common = set(sentences1) & set(sentences2)
    only_in_1 = set(sentences1) - set(sentences2)
    only_in_2 = set(sentences2) - set(sentences1)
    
    # Calculate similarity percentage
    total_sentences = len(set(sentences1) | set(sentences2))
    similarity = len(common) / total_sentences if total_sentences > 0 else 0
    
    # Format differences
    differences = []
    if only_in_1:
        differences.append({
            "type": "removed",
            "sentences": list(only_in_1)
        })
    if only_in_2:
        differences.append({
            "type": "added",
            "sentences": list(only_in_2)
        })
    
    return {
        "similarity_percentage": round(similarity * 100, 2),
        "differences": differences,
        "common_points": list(common)
    } 